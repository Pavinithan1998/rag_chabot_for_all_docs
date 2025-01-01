
import fitz  
from io import BytesIO
from docx import Document
from pptx import Presentation
from llm_actions import get_image_description


def process_pdf(file_bytes: bytes, filename: str):
    """Process PDF file to extract text and image descriptions, integrating them into the text content."""
    text_content = ""
    pdf_document = fitz.open("pdf", file_bytes)
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        text = page.get_text()
        text_content += f"Page {page_num + 1}:\n{text}\n"

        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            try:
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = f"page_{page_num + 1}_img_{img_index + 1}.{image_ext}"
                
                description = get_image_description(image_bytes)
                
                text_content += f"\n[Image: {image_filename}]\nDescription: {description}\n"
            
            except Exception as e:
                text_content += f"\n[Image Extraction Failed: Page {page_num + 1}, Image {img_index + 1}]\nError: {e}\n"
    # save the text content to a txt file
    file_path = f"./docs/{filename}.txt"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text_content)
    return text_content

def process_docx(file_bytes: bytes, filename: str):
    """Process DOCX file using fitz to extract text and image descriptions."""
    text_content = ""
    docx_pdf = fitz.open("docx", file_bytes)
    
    for page_num in range(docx_pdf.page_count):
        page = docx_pdf.load_page(page_num)
        text = page.get_text()
        text_content += f"Page {page_num + 1}:\n{text}\n"
        
        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            try:
                base_image = docx_pdf.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = f"page_{page_num + 1}_img_{img_index + 1}.{image_ext}"
                
                description = get_image_description(image_bytes)
                
                text_content += f"\n[Image: {image_filename}]\nDescription: {description}\n"
            except Exception as e:
                text_content += f"\n[Image Extraction Failed: Page {page_num + 1}, Image {img_index + 1}]\nError: {e}\n"
    # save the text content to a txt file
    file_path = f"./docs/{filename}.txt"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text_content)
    return text_content

def process_txt(file_bytes: bytes, filename: str):
    """Process TXT file to extract plain text and store in Pinecone vector store."""
    text_content = file_bytes.decode("utf-8")
    # save the text content to a txt file
    file_path = f"./docs/{filename}.txt"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text_content)
    return text_content


def process_ppt(file_bytes: bytes, filename: str):
    """Process PPT file to extract text and image descriptions, integrating them into the text content."""
    text_content = ""
    image_descriptions = {}
    ppt = Presentation(BytesIO(file_bytes))
    
    for slide_num, slide in enumerate(ppt.slides):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + "\n"
        text_content += f"Slide {slide_num + 1}:\n{slide_text}\n"
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 is the shape type for images in PPTX
                image = shape.image.blob
                image_filename = f"slide_{slide_num + 1}_img_{len(image_descriptions) + 1}.jpg"
                description = get_image_description(image)
                image_descriptions[image_filename] = description
                text_content += f"\n[Image: {image_filename}]\nDescription: {description}\n"    
    # save the text content to a txt file
    file_path = f"./docs/{filename}.txt"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text_content)
    return text_content